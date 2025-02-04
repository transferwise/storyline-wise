from pptx import Presentation
from pptx.util import Inches, Pt
from markdown import markdown
from bs4 import BeautifulSoup
import os


def markdown_to_slide(md_section: str, filename: str):
    """
    Convert a section of markdown into a PowerPoint slide.

    Parameters:
        md_section (str): The Markdown content to be added to the slide.
        filename (str): The PowerPoint file to save or append to.
    """
    # Load or create a PowerPoint presentation
    if os.path.exists(filename):
        prs = Presentation(filename)
    else:
        prs = Presentation()

    # Create a new slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content layout

    # Parse the Markdown content
    html = markdown(md_section, extensions=["extra", "nl2br", "sane_lists"])
    soup = BeautifulSoup(html, "html.parser")

    # Extract the first heading for the slide title
    title = slide.shapes.title
    heading = soup.find(["h1", "h2", "h3", "h4", "h5", "h6"])
    if heading:
        title.text = heading.get_text()
        heading.extract()  # Remove the title from the content
    else:
        title.text = "Slide"

    # Extract the remaining content for the slide body
    content = slide.placeholders[1]
    content.text = ""  # Clear any default text

    for element in soup:
        if element.name == "p":  # Process paragraph tags as plain text
            para = content.text_frame.add_paragraph()
            para.text = element.get_text()
            para.space_after = Pt(10)  # Add spacing between paragraphs
            para.space_before = Pt(10)
            para.level = 0  # No bullet point
        elif element.name == "ul":  # Process unordered lists
            for li in element.find_all("li"):
                para = content.text_frame.add_paragraph()
                para.text = f"{li.get_text()}"
                para.level = 1  # Bullet point level 1
        elif element.name == "ol":  # Process ordered lists (numbered)
            for li in element.find_all("li"):
                para = content.text_frame.add_paragraph()
                para.text = li.get_text()
                para.level = 1  # Same level as bullet points
                # Set paragraph to use numbered list formatting
                para._element.get_or_add_pPr().set("numPr", "")

    # Save the PowerPoint file
    prs.save(filename)
    print(f"Slide added to {filename}")


if __name__ == "__main__":
    # All indented lines are interpreted as code, so remove indents
    md_content = """# Markdown Slide Title
This is a paragraph in the slide content.

- First bullet point
- Second bullet point

1. First numbered item
2. Second numbered item
"""

    markdown_to_slide(md_content, "example_presentation.pptx")
    print("Done!")
